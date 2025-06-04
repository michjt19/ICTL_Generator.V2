from pptx import Presentation

def generate_ppt(task, filepath):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = task["title"]
    slide.placeholders[1].text = f"Task Code: {task['code']}"

    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Condition & Standard"
    content = slide.placeholders[1]
    content.text = f"Condition: {task['condition']}\nStandard: {task['standard']}"

    if "steps" in task:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Performance Steps"
        body = slide.placeholders[1]
        for step in task["steps"]:
            body.text += f"\n• {step}"

    prs.save(filepath)
